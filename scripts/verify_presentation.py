import os
from pptx import Presentation

pptx_path = 'AI_Text_Detector_Presentation.pptx'
prs = Presentation(pptx_path)

print(f"File size: {os.path.getsize(pptx_path)} bytes")
print(f"Total slides count: {len(prs.slides)}")

picture_count = 0
for idx, slide in enumerate(prs.slides):
    texts = [s.text_frame.text.strip().replace('\n', ' ') for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    pics = [s for s in slide.shapes if s.shape_type == 13 or hasattr(s, 'image')]
    picture_count += len(pics)
    first_text = texts[0][:60] if texts else "No Text"
    print(f"Slide {idx+1:02d}: {len(texts)} text boxes, {len(pics)} pictures | Header: {first_text}")

print(f"Total pictures embedded across presentation: {picture_count}")
