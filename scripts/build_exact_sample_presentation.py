import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt

def update_slide_title(slide, new_title):
    for shape in slide.shapes:
        if shape.shape_type == 6: # Group
            if hasattr(shape, 'shapes'):
                for s in shape.shapes:
                    if s.has_text_frame and (s.name == "TextBox 4" or s.name == "TextBox 9"):
                        tf = s.text_frame
                        if tf.paragraphs:
                            p = tf.paragraphs[0]
                            r0 = p.runs[0] if p.runs else None
                            f_name = r0.font.name if r0 and r0.font else None
                            f_size = r0.font.size.pt if r0 and r0.font and r0.font.size else None
                            f_bold = r0.font.bold if r0 and r0.font else None
                            f_color = r0.font.color.rgb if r0 and r0.font and hasattr(r0.font.color, 'rgb') else None

                            p.text = new_title

                            if p.runs:
                                r = p.runs[0]
                                if f_name: r.font.name = f_name
                                if f_size: r.font.size = Pt(f_size)
                                if f_bold is not None: r.font.bold = f_bold
                                if f_color: r.font.color.rgb = f_color
                        return
        elif shape.has_text_frame and shape.top < Inches(1.5):
            tf = shape.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]
                if p.text.strip():
                    p.text = new_title
                    return

def update_text_frame_bullets(tf, paragraphs_data):
    while len(tf.paragraphs) < len(paragraphs_data):
        tf.add_paragraph()

    for i, p_data in enumerate(paragraphs_data):
        if i >= len(tf.paragraphs):
            break
        text = p_data[0]
        level = p_data[1] if len(p_data) > 1 else 0

        p = tf.paragraphs[i]
        
        orig_name = None
        orig_size = None
        orig_bold = None
        orig_color = None
        
        if p.runs:
            r0 = p.runs[0]
            if r0.font:
                orig_name = r0.font.name
                orig_size = r0.font.size.pt if r0.font.size else None
                orig_bold = r0.font.bold
                if hasattr(r0.font.color, 'rgb'):
                    orig_color = r0.font.color.rgb

        p.text = text
        p.level = level

        for run in p.runs:
            if orig_name: run.font.name = orig_name
            if orig_size: run.font.size = Pt(orig_size)
            if orig_bold is not None: run.font.bold = orig_bold
            if orig_color: run.font.color.rgb = orig_color

def replace_image_in_slide(slide, img_path, left=Inches(1.5), top=Inches(1.5), width=Inches(7.6), height=Inches(5.2)):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)

def build_presentation():
    prs = Presentation('Sample.pptx')
    
    # SLIDE 1: Title Slide
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.shape_type == 6 and shape.name == "Group 2":
            for g_s in shape.shapes:
                if g_s.has_text_frame and g_s.name == "TextBox 4":
                    tf = g_s.text_frame
                    for p in tf.paragraphs:
                        t = p.text.strip()
                        if "Web Application Vulnerability Scanner" in t:
                            p.text = "AI Text Detector"
                        elif "Shashank Singh" in t:
                            p.text = "AI Text Detector Team"
                        elif "Bachelor in Computer Applications" in t:
                            p.text = "A Final Project Presentation Submitted in partial fulfillment of the requirement of Bachelor of Science in Computer Science and Information Technology (B.Sc. CSIT) 8th semester"

    # SLIDE 2: Introduction to project
    update_slide_title(prs.slides[1], "Introduction to project")
    for shape in prs.slides[1].shapes:
        if shape.name == "TextBox 5":
            update_text_frame_bullets(shape.text_frame, [
                ("AI Text Detector ('Sentinel') is a client-side web application designed to evaluate document authenticity and detect synthetic AI-generated text.", 1),
                ("With the rapid rise of Large Language Models (ChatGPT, Claude, Gemini), automated AI text detection has become essential for academic integrity and publishing.", 1),
                ("Traditional cloud detectors introduce high privacy risks by transmitting sensitive documents over third-party networks, alongside high API latency.", 1),
                ("The AI Text Detector runs entirely in the browser sandbox, performing instant stylometric evaluation, statistical NLP modeling, and ML classification.", 1),
                ("Built with Vanilla JavaScript, HTML5, CSS3, MathJax, and Scikit-Learn, the tool provides privacy-preserving detection with interactive visual heatmaps.", 1)
            ])

    # SLIDE 3: Objective of the Project
    update_slide_title(prs.slides[2], "Objective of the Project")
    for shape in prs.slides[2].shapes:
        if shape.name == "TextBox 5":
            update_text_frame_bullets(shape.text_frame, [
                ("To design and implement a 100% client-side AI Text Detector that automatically identifies machine-generated synthetic text directly within the browser.", 1),
                ("To integrate statistical NLP (Perplexity & Burstiness), KNN, Logistic Regression, Character Trigrams, POS Syntax profiling, and Heuristic Grammar calibration into an explainable visual dashboard.", 1)
            ])

    # SLIDE 4: Scope of Project
    update_slide_title(prs.slides[3], "Scope of Project")
    for shape in prs.slides[3].shapes:
        if shape.name == "TextBox 5":
            update_text_frame_bullets(shape.text_frame, [
                ("Automatic detection of synthetic text generated by modern LLMs using multi-dimensional stylometric features, sentence-level perplexity heatmaps, and LaTeX algorithmic transparency panels.", 1),
                ("Real-time client-side document parsing (.TXT, .PDF, .DOCX), interactive SVG chart rendering (Gauge dial, 2D KNN scatter plot, 5-axis radar web), and zero server-side data retention.", 1)
            ])

    # SLIDE 5: Functional Requirement (Use Case Diagram)
    update_slide_title(prs.slides[4], "Functional Requirement")
    replace_image_in_slide(prs.slides[4], 'extracted_images/image_04.png', Inches(2.2), Inches(1.3), Inches(6.2), Inches(5.0))

    # SLIDE 6: Feasibility Study - Technical Hardware
    update_slide_title(prs.slides[5], "Feasibility Study")
    for shape in prs.slides[5].shapes:
        if shape.name == "TextBox 5":
            update_text_frame_bullets(shape.text_frame, [
                ("Technical", 1),
                ("Hardware Requirements:", 2),
                ("RAM Size: A minimum of 4 GB RAM is recommended to ensure smooth client-side execution of JavaScript tokenization, MathJax 3 rendering, and interactive SVG graphics.", 3),
                ("Processor Size: A standard dual-core CPU with at least 2.0 GHz clock speed efficiently performs text tokenization, matrix vector dot-products, and ensemble scoring in under 500ms.", 3),
                ("Storage Space: Less than 50 MB of local storage is required for static HTML/JS web assets, CSS styles, and pre-computed machine learning model vectors.", 3),
                ("Network Equipment: No active internet connection is required for core detection, ensuring full offline functionality and complete data privacy.", 3)
            ])

    # SLIDE 7: Software Requirements
    for shape in prs.slides[6].shapes:
        if shape.name == "TextBox 2":
            update_text_frame_bullets(shape.text_frame, [
                ("Software Requirements:", 2),
                ("Operating System: Compatible with Windows 10/11, Linux, macOS, and ChromeOS across any modern Web browser runtime environment.", 3),
                ("Frontend Technologies: Built using HTML5, CSS3 Glassmorphism styles, Vanilla ES6+ JavaScript, PDF.js for PDF parsing, and Mammoth.js for Word document parsing.", 3),
                ("Machine Learning & Typesetting Tools: Uses Scikit-Learn and NLTK for offline model generation, and MathJax 3 for dynamic LaTeX algorithmic formula typesetting.", 3)
            ])

    # SLIDE 8: Operational Feasibility
    for shape in prs.slides[7].shapes:
        if shape.name == "TextBox 2":
            update_text_frame_bullets(shape.text_frame, [
                ("Operational", 1),
                ("The AI Text Detector is designed with an intuitive drag-and-drop dashboard interface that enables educators, students, and researchers to conduct text authenticity checks effortlessly.", 2),
                ("The system automates stylometric feature extraction, providing instant visual feedback via semi-circular gauge dials, 2D KNN scatter plots, 5-axis radar webs, and sentence-level heatmaps.", 2)
            ])

    # SLIDE 9: Economical Feasibility
    for shape in prs.slides[8].shapes:
        if shape.name == "TextBox 2":
            update_text_frame_bullets(shape.text_frame, [
                ("Economical", 1),
                ("The AI Text Detector is economically feasible because it is developed entirely using open-source technologies (HTML5, Vanilla JS, CSS3, MathJax, Scikit-Learn), eliminating all cloud hosting and API subscription costs.", 2),
                ("Development utilities and client-side libraries are freely available under MIT/Apache licenses, resulting in zero operational server overhead.", 2),
                ("By executing all computation client-side in browser memory, the system provides unlimited text analysis without incurring third-party API token charges.", 2)
            ])

    # SLIDE 10: Schedule (Gantt Chart)
    update_slide_title(prs.slides[9], "Schedule")
    replace_image_in_slide(prs.slides[9], 'extracted_images/image_03.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.0))

    # SLIDE 11: Methodology
    update_slide_title(prs.slides[10], "Methodology")

    # SLIDE 12: System Design
    update_slide_title(prs.slides[11], "System Design")

    # SLIDE 13: Steps in Stylometric Algorithm
    for shape in prs.slides[12].shapes:
        if shape.name == "TextBox 2":
            update_text_frame_bullets(shape.text_frame, [
                ("Steps in Stylometric & Machine Learning Algorithms for AI Text Detector:", 1),
                ("Data Collection & Pre-processing: Document text is parsed (.TXT, .PDF, .DOCX), tokenized into sentences/words, and stripped of invalid characters.", 2),
                ("Feature Vector Extraction: Computes 5D stylometric metrics including Type-Token Ratio (TTR), Sentence Burstiness, AI Vocabulary Density, Readability Index, and Punctuation Density.", 2),
                ("Multi-Engine Classification & Voting: Evaluates text across 6 independent base models (Logistic Regression, Character Trigrams, KNN, POS Syntax, Cosine Similarity, Perplexity).", 2),
                ("Grammar Factor Calibration: Applies Heuristic Grammar & Typography scoring to calibrate the ensemble output against human typing errors.", 2),
                ("Alert Threshold System (Probability Insight):", 2),
                ("Red Alert: AI Probability > 60% indicates synthetic generated text (Extremely High / High AI Probability).", 3),
                ("Green Alert: AI Probability < 35% indicates authentic human writing (Negligible / Human-Like Text).", 3)
            ])

    # SLIDES 14 - 19: UML Diagrams
    update_slide_title(prs.slides[13], "Class Diagram")
    replace_image_in_slide(prs.slides[13], 'extracted_images/image_06.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))
    
    update_slide_title(prs.slides[14], "Sequence Diagram")
    replace_image_in_slide(prs.slides[14], 'extracted_images/image_08.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))
    
    update_slide_title(prs.slides[15], "State Diagram")
    replace_image_in_slide(prs.slides[15], 'extracted_images/image_07.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))
    
    update_slide_title(prs.slides[16], "Activity Diagram")
    replace_image_in_slide(prs.slides[16], 'extracted_images/image_09.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))
    
    update_slide_title(prs.slides[17], "Component Diagram")
    replace_image_in_slide(prs.slides[17], 'extracted_images/image_10.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))
    
    update_slide_title(prs.slides[18], "Deployment Diagram")
    replace_image_in_slide(prs.slides[18], 'extracted_images/image_11.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))

    # SLIDE 20: Tools Used
    update_slide_title(prs.slides[19], "Tools Used")
    for shape in prs.slides[19].shapes:
        if shape.name == "TextBox 5":
            update_text_frame_bullets(shape.text_frame, [
                ("Front-End Tools:", 0),
                ("HTML5 / CSS3 Glassmorphism", 1),
                ("Vanilla JavaScript (ES6+)", 1),
                ("MathJax 3 & PDF.js / Mammoth.js", 1),
                ("Backend & ML Tools:", 0),
                ("Python 3.10", 1),
                ("Scikit-Learn & NLTK", 1),
                ("NLPAug", 1)
            ])

    # SLIDE 21: Modules List
    update_slide_title(prs.slides[20], "Modules")
    for shape in prs.slides[20].shapes:
        if shape.name == "TextBox 5":
            update_text_frame_bullets(shape.text_frame, [
                ("FileParser Module (.TXT, .PDF, .DOCX Ingestion)", 1),
                ("Statistical Perplexity & Burstiness Engine", 1),
                ("K-Nearest Neighbors (KNN) Feature Vector Classifier", 1),
                ("Stylometric Logistic Regression Classifier", 1),
                ("Character Trigram Cosine Similarity Engine", 1),
                ("Part-of-Speech (POS) Syntax Ratios Classifier", 1),
                ("Heuristic Grammar & Typography Calibrator", 1),
                ("SVG Semi-Circular Gauge Renderer Module", 1),
                ("2D KNN Feature Space Scatter Plot Module", 1),
                ("5-Axis Syntactic Radar Web Module", 1),
                ("Sentence-Level Perplexity Heatmap Engine", 1),
                ("MathJax Algorithmic Resolution Panel Module", 1),
                ("Ensemble Scoring & Event Controller Sandbox", 1)
            ])

    # SLIDES 22 - 34: 13 Individual Modules
    modules_data = [
        ("FileParser Module",
         "The FileParser Module enables multi-format document ingestion, handling plain text files (.TXT, .MD), PDF documents via Mozilla's PDF.js canvas parser, and Word documents via Mammoth.js XML converter.",
         "It validates the extracted text length against a 10-word minimum threshold required for stylometric confidence and caches normalized sentence/word tokens in browser memory."),
         
        ("Statistical Perplexity & Burstiness Engine",
         "The Perplexity & Burstiness Module evaluates local n-gram likelihoods and sentence length variance across the document text.",
         "It identifies uniform sentence structures typical of AI models and flags unexpected vocabulary spikes indicative of authentic human writing."),
         
        ("KNN Feature Vector Classifier Module",
         "The K-Nearest Neighbors (KNN) Module maps document perplexity and sentence variance into a 2D normalized feature coordinate space.",
         "It computes Euclidean distances to pre-loaded training centroids (Human vs. AI text clusters) and passes 2D coordinate points to the visual scatter plot renderer."),
         
        ("Stylometric Logistic Regression Model",
         "The Stylometric Logistic Regression Module evaluates a 5-dimensional feature vector comprising Type-Token Ratio, Burstiness, AI Vocabulary Density, Readability Index, and Punctuation Density.",
         "It applies pre-trained logistic regression weights trained on 1,000 document samples to output a high-precision binary AI probability score."),
         
        ("Character Trigram Cosine Engine",
         "The Character Trigram Module extracts overlapping 3-character sub-word tokens across the text to capture morpho-syntactic writing signatures.",
         "It measures cosine vector similarity against AI and Human prototype vectors, rendering the model resilient against word-level synonym substitution evasions."),
         
        ("POS Syntax Ratios Classifier Module",
         "The POS Syntax Ratios Module profiles grammatical structure across Noun, Verb, Adjective, Adverb, and Function Word distributions.",
         "It compares syntactic distribution ratios against baseline prototype vectors and supplies 5-axis coordinate data to the interactive SVG Radar Web visualizer."),
         
        ("Heuristic Grammar & Typography Calibrator",
         "The Heuristic Grammar & Typography Module evaluates text for grammatical perfection, punctuation patterns, subject-verb agreement, and typing flaws.",
         "It computes a Grammar Perfection Score that acts as a calibration scaling factor, rewarding authentic human typing errors and penalizing synthetic perfection."),
         
        ("SVG Semi-Circular Gauge Renderer",
         "The SVG Gauge Module dynamically renders a semi-circular vector dial in the dashboard UI using trigonometric arc trajectory calculations.",
         "It interpolates color transitions from Sage Green (Human) to Terracotta Red (AI) with smooth 800ms needle rotation animations."),
         
        ("2D KNN Feature Scatter Plot Module",
         "The 2D Scatter Plot Module draws an interactive SVG coordinate graph displaying reference training dots alongside the active document coordinate marker.",
         "It provides spatial visual proof of proximity to Human vs. AI reference clusters, enhancing model explainability."),
         
        ("5-Axis Syntactic Radar Web Module",
         "The 5-Axis Radar Web Module renders an interactive SVG spider graph representing document syntax dimensions.",
         "It overlays the active document footprint against Human and AI benchmark polygons to highlight structural syntactic deviations."),
         
        ("Sentence-Level Perplexity Heatmap Engine",
         "The Sentence Heatmap Engine tokenizes text into individual sentences and computes localized perplexity scores for each sentence independently.",
         "It renders color-coded background highlights directly in the editor sandbox with hover tooltips displaying numerical perplexity metrics."),
         
        ("MathJax Algorithmic Resolution Panel",
         "The MathJax Resolution Panel dynamically formats step-by-step feature vector coordinates, weights, and equations into LaTeX mathematical notation.",
         "It invokes MathJax 3 to render publication-grade math formulas, providing complete transparency into the ensemble classification calculation."),
         
        ("Ensemble Scoring & Event Controller",
         "The Ensemble Scoring Orchestrator coordinates parallel execution across all 6 base classification engines and integrates outputs into a weighted ensemble score.",
         "It applies grammar calibration multipliers, maps final scores to probability risk tiers, and manages DOM UI state updates seamlessly.")
    ]

    for m_idx, (m_title, p0_text, p1_text) in enumerate(modules_data):
        slide_idx = 21 + m_idx # Slides 22 to 34
        slide = prs.slides[slide_idx]
        
        update_slide_title(slide, m_title)
        
        for shape in slide.shapes:
            if shape.name == "TextBox 10":
                update_text_frame_bullets(shape.text_frame, [
                    (p0_text, 1),
                    (p1_text, 1)
                ])

    # SLIDE 35: Testing
    update_slide_title(prs.slides[34], "Testing")
    for shape in prs.slides[34].shapes:
        if shape.name == "TextBox 5":
            update_text_frame_bullets(shape.text_frame, [
                ("Unit Testing", 0),
                ("System Testing", 0)
            ])

    # SLIDES 36 - 38: Results & Confusion Matrix
    update_slide_title(prs.slides[35], "Classifier Performance Analysis")
    replace_image_in_slide(prs.slides[35], 'extracted_images/image_13.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))
    
    update_slide_title(prs.slides[36], "Test Cases & Score Distribution")
    replace_image_in_slide(prs.slides[36], 'extracted_images/image_12.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))
    
    update_slide_title(prs.slides[37], "Confusion Matrix & Evaluation")
    replace_image_in_slide(prs.slides[37], 'extracted_images/image_16.png', Inches(1.5), Inches(1.5), Inches(7.6), Inches(5.2))

    # SLIDE 39: Conclusion
    update_slide_title(prs.slides[38], "Conclusion")
    for shape in prs.slides[38].shapes:
        if shape.name == "TextBox 5":
            update_text_frame_bullets(shape.text_frame, [
                ("The AI Text Detector system automatically detects synthetic machine-generated text using 6 complementary statistical NLP and machine learning engines.", 1),
                ("Built using Vanilla JavaScript, HTML5, CSS3, and MathJax, the system operates 100% client-side with zero data privacy risks and zero server hosting costs.", 1),
                ("The system efficiently performs real-time feature extraction, calculates risk probabilities (94.2% accuracy), and renders interactive SVG charts and sentence heatmaps.", 1),
                ("Overall, the project helps educators, students, and researchers verify document authenticity while maintaining complete data privacy and transparency.", 1)
            ])

    # SLIDE 40: DEMO
    update_slide_title(prs.slides[39], "DEMO")

    output_path = "AI_Text_Detector_Presentation.pptx"
    prs.save(output_path)
    print(f"Direct modification saved cleanly to {output_path}")

if __name__ == "__main__":
    build_presentation()
