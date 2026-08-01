from pptx import Presentation

prs = Presentation('Sample.pptx')
slide22 = prs.slides[21]
for shape in slide22.shapes:
    print(f"Shape: name='{shape.name}', type={shape.shape_type}")
    if shape.shape_type == 6: # Group
        for s in shape.shapes:
            print(f"  Sub-shape: name='{s.name}', type={s.shape_type}")
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    print(f"    Text: '{p.text}'")
