from pptx import Presentation

prs = Presentation('AI_Text_Detector_Presentation.pptx')
slide22 = prs.slides[21]
for i, s in enumerate(slide22.shapes):
    if s.has_text_frame:
        print(f"Shape {i} ({s.name}): '{s.text_frame.text.strip()}'")
    elif s.shape_type == 6: # Group
        for sub in s.shapes:
            if sub.has_text_frame:
                print(f"Group sub-shape ({sub.name}): '{sub.text_frame.text.strip()}'")
