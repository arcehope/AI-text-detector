import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation('Sample.pptx')
with open('sample_structure_dump.txt', 'w', encoding='utf-8') as f:
    for idx, slide in enumerate(prs.slides):
        f.write(f"*** SLIDE {idx+1:02d} ***\n")
        for s_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                f.write(f"  Shape {s_idx} ({shape.name}):\n")
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    text = p.text.strip()
                    if text:
                        f.write(f"    P{p_idx} (lvl={p.level}): '{text}'\n")
            elif shape.shape_type == 13: # Picture
                f.write(f"  Shape {s_idx} [PICTURE]: name={shape.name}, left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}\n")
            elif shape.shape_type == 6: # Group
                f.write(f"  Shape {s_idx} [GROUP]: name={shape.name}\n")
                if hasattr(shape, 'shapes'):
                    for g_s in shape.shapes:
                        if g_s.has_text_frame:
                            for p_idx, p in enumerate(g_s.text_frame.paragraphs):
                                if p.text.strip():
                                    f.write(f"      Sub-shape {g_s.name} P{p_idx} (lvl={p.level}): '{p.text.strip()}'\n")

print("Dumped to sample_structure_dump.txt successfully.")
