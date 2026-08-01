from pptx import Presentation

prs = Presentation('AI_Text_Detector_Presentation.pptx')
for idx in range(21, 34):
    slide = prs.slides[idx]
    header_text = ''
    body_text = ''
    for shape in slide.shapes:
        if shape.shape_type == 6: # Group
            if hasattr(shape, 'shapes'):
                for g_s in shape.shapes:
                    if g_s.has_text_frame:
                        header_text = g_s.text_frame.text.strip()
        elif shape.has_text_frame:
            body_text = shape.text_frame.text.strip()[:60]
    print(f"Slide {idx+1}: Header='{header_text}' | Body='{body_text}'")
