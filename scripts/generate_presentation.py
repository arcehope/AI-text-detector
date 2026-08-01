import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank layout

    # Color Palette Definitions
    DARK_BG = RGBColor(15, 23, 42)        # Slate 900
    SLATE_CARD = RGBColor(30, 41, 59)     # Slate 800
    LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
    CARD_BG = RGBColor(255, 255, 255)    # Pure White
    CARD_BORDER = RGBColor(226, 232, 240)# Slate 200
    
    PRIMARY_BLUE = RGBColor(37, 99, 235) # Blue 600
    DARK_BLUE = RGBColor(30, 58, 138)    # Blue 900
    TEAL_ACCENT = RGBColor(13, 148, 136) # Teal 600
    CYAN_ACCENT = RGBColor(6, 182, 212)  # Cyan 500
    INDIGO_ACCENT = RGBColor(79, 70, 229)# Indigo 600
    EMERALD_ACCENT = RGBColor(16, 185, 129)# Emerald 500
    AMBER_ACCENT = RGBColor(217, 119, 6) # Amber 600
    ROSE_ACCENT = RGBColor(225, 29, 72)  # Rose 600
    
    TEXT_DARK = RGBColor(30, 41, 59)     # Slate 800
    TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
    TEXT_LIGHT = RGBColor(241, 245, 249) # Slate 100

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_badge="SENTINEL AI TEXT DETECTOR", dark=False):
        # Header Container Shape
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Category Badge
        p0 = tf.paragraphs[0]
        p0.text = category_badge.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.name = 'Calibri'
        p0.font.color.rgb = CYAN_ACCENT if dark else PRIMARY_BLUE
        p0.space_after = Pt(2)
        
        # Title
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.name = 'Calibri'
        p1.font.color.rgb = TEXT_LIGHT if dark else DARK_BLUE

        # Top Right Line Decorator
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = PRIMARY_BLUE if not dark else CYAN_ACCENT
        line.line.fill.background()

    def add_footer(slide, current_page, total_pages=40, dark=False):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.4))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = "Sentinel: AI Text Detector Project Defense | Tribhuvan University"
        p.font.size = Pt(10)
        p.font.name = 'Calibri'
        p.font.color.rgb = TEXT_MUTED if not dark else RGBColor(148, 163, 184)
        
        # Page Number Right Aligned
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.RIGHT
        p2.text = f"{current_page} / {total_pages}"
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.name = 'Calibri'
        p2.font.color.rgb = PRIMARY_BLUE if not dark else CYAN_ACCENT

    def add_card(slide, left, top, width, height, title, items, badge="", bg_color=CARD_BG, border_color=CARD_BORDER, dark_text=True):
        # Card background rectangle
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        
        # Inner text frame
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        if badge:
            p_badge = tf.paragraphs[0]
            p_badge.text = badge.upper()
            p_badge.font.size = Pt(9)
            p_badge.font.bold = True
            p_badge.font.name = 'Calibri'
            p_badge.font.color.rgb = PRIMARY_BLUE if dark_text else CYAN_ACCENT
            p_badge.space_after = Pt(2)
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]
            
        p_title.text = title
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.name = 'Calibri'
        p_title.font.color.rgb = TEXT_DARK if dark_text else TEXT_LIGHT
        p_title.space_after = Pt(8)
        
        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = f"•  {item}"
            p_item.font.size = Pt(12)
            p_item.font.name = 'Calibri'
            p_item.font.color.rgb = TEXT_DARK if dark_text else RGBColor(226, 232, 240)
            p_item.space_after = Pt(4)

    def add_image_slide(slide, img_path, left, top, width, height, caption=""):
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            if caption:
                cap_box = slide.shapes.add_textbox(left, top + height + Inches(0.05), width, Inches(0.35))
                tf = cap_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.text = caption
                p.font.size = Pt(11)
                p.font.italic = True
                p.font.name = 'Calibri'
                p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (Dark Executive Theme)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, DARK_BG)
    
    # Title Box
    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "SENTINEL: AI TEXT DETECTOR"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(10)
    
    p2 = tf1.add_paragraph()
    p2.text = "A Multidimensional Stylometric & Machine Learning System for Client-Side Synthetic Content Identification"
    p2.font.size = Pt(18)
    p2.font.name = 'Calibri'
    p2.font.color.rgb = CYAN_ACCENT
    p2.space_after = Pt(36)
    
    p3 = tf1.add_paragraph()
    p3.text = "Project Defense Presentation | B.Sc. CSIT / Bachelor's Final Project"
    p3.font.size = Pt(14)
    p3.font.name = 'Calibri'
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_after = Pt(6)
    
    p4 = tf1.add_paragraph()
    p4.text = "Department of Computer Science & Information Technology\nTribhuvan University, Nepal"
    p4.font.size = Pt(13)
    p4.font.name = 'Calibri'
    p4.font.color.rgb = TEXT_MUTED

    # Decorative bottom accent line
    dec = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(6.0), Inches(4.0), Inches(0.06))
    dec.fill.solid()
    dec.fill.fore_color.rgb = PRIMARY_BLUE
    dec.line.fill.background()

    # Add TU Logo image if exists
    if os.path.exists('extracted_images/image_01.png'):
        slide1.shapes.add_picture('extracted_images/image_01.png', Inches(10.2), Inches(1.8), width=Inches(2.0))

    # =========================================================================
    # SLIDE 2: INTRODUCTION & BACKGROUND
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, LIGHT_BG)
    add_header(slide2, "Introduction & Background Context", "CHAPTER 1: INTRODUCTION")
    add_footer(slide2, 2)
    
    add_card(slide2, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Rise of Synthetic AI Text",
             ["Proliferation of Large Language Models (ChatGPT, Claude, Gemini, Llama).",
              "Mass generation of human-like essays, articles, research papers, and reports.",
              "Blurring boundaries between authentic human authorship and synthetic output."],
             badge="Context & Proliferation")
             
    add_card(slide2, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Key Domain Challenges",
             ["Academic integrity risks (unattributed AI-assisted plagiarism).",
              "Journalism & publishing authenticity concerns.",
              "Misinformation scaling and automated content farm generation.",
              "Legal & copyright verification requirements."],
             badge="Industry Impact")
             
    add_card(slide2, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "The Sentinel Solution",
             ["Client-side AI detection running entirely in the browser.",
              "Combines 6 independent NLP & Machine Learning engines.",
              "Zero data transmission — guarantees 100% user document privacy.",
              "Provides sentence-level heatmaps and interactive chart diagnostics."],
             badge="Proposed Paradigm")

    # =========================================================================
    # SLIDE 3: PROBLEM STATEMENT
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, LIGHT_BG)
    add_header(slide3, "Problem Statement & Existing Limitations", "PROBLEM ANALYSIS")
    add_footer(slide3, 3)

    add_card(slide3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.4),
             "Privacy & Security Vulnerabilities",
             ["Existing commercial detectors (Turnitin AI, GPTZero, CopyLeaks) process documents on remote cloud servers.",
              "Exposes sensitive academic drafts, trade secrets, and personal documents to third-party retention and leakage."],
             badge="01. Cloud Dependency")

    add_card(slide3, Inches(6.8), Inches(1.6), Inches(5.6), Inches(2.4),
             "Black-Box Predictions & Low Explainability",
             ["Most detectors return a single opaque percentage without justification.",
              "Users cannot inspect which specific sentences or stylometric patterns triggered the classification decision."],
             badge="02. Opacity")

    add_card(slide3, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.4),
             "High False Positive Rates on ESL Writers",
             ["Deep neural detectors misclassify non-native English writing due to repetitive syntax structures.",
              "Creates unfair academic penalties and mistrust among students and educators."],
             badge="03. Bias & Misclassification")

    add_card(slide3, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.4),
             "High Latency & Server Cost Overhead",
             ["Cloud API calls incur high subscription costs and multi-second network latency.",
              "Unsuitable for offline evaluation, confidential institutions, or instant interactive analysis."],
             badge="04. Cost & Latency")

    # =========================================================================
    # SLIDE 4: PROJECT OBJECTIVES
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, LIGHT_BG)
    add_header(slide4, "Project Objectives & Technical Scope", "PROJECT GOALS")
    add_footer(slide4, 4)

    add_card(slide4, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Primary Objective",
             ["Design and implement a 100% client-side AI text detection system ('Sentinel').",
              "Execute all feature extraction, ML classification, and UI rendering locally in the browser sandbox.",
              "Achieve high detection accuracy without sending a single byte of data over the internet."],
             badge="Core Objective")

    add_card(slide4, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Specific Technical Goals",
             ["Implement statistical NLP engines (N-gram Perplexity & Burstiness).",
              "Train and integrate KNN and Stylometric Logistic Regression classifiers.",
              "Build character trigram and POS syntax vector space similarity engines.",
              "Calibrate predictions with a Heuristic Grammar & Typography analyzer."],
             badge="Algorithmic Goals")

    add_card(slide4, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "UI & Explainability Goals",
             ["Render real-time interactive SVG charts (Gauge, 2D KNN Scatter, 5-Axis Radar).",
              "Highlight sentence-level perplexity heatmaps directly in the document editor.",
              "Typeset exact mathematical calculations via MathJax 3 LaTeX integration."],
             badge="Explainability Goals")

    # =========================================================================
    # SLIDE 5: DYNAMIC MODELING - USE CASE DIAGRAM
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, LIGHT_BG)
    add_header(slide5, "System Use Case Modeling", "CHAPTER 3: SYSTEM ANALYSIS")
    add_footer(slide5, 5)

    add_card(slide5, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Use Case Specification",
             ["Actor: End User (Student, Educator, Researcher).",
              "UC-1: Document or Text Input (Paste text or drag-and-drop .txt, .pdf, .docx).",
              "UC-2: Sample Text Loader (Load pre-set AI/Human benchmarks).",
              "UC-3: Run Analysis (Triggers 6 parallel base engines).",
              "UC-4: View Visual Diagnostics (Gauge dial, KNN scatter plot, 5-axis radar web).",
              "UC-5: Inspect Perplexity Heatmap (Sentence-by-sentence background highlights).",
              "UC-6: Review LaTeX Resolution Panel (Exact feature vector breakdown)."],
             badge="Use Case Breakdown")

    add_image_slide(slide5, 'extracted_images/image_04.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure 3.1: Use Case Diagram for AI Text Detector")

    # =========================================================================
    # SLIDE 6: TECHNICAL FEASIBILITY - HARDWARE REQUIREMENTS
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, LIGHT_BG)
    add_header(slide6, "Technical Feasibility: Hardware Requirements", "FEASIBILITY ANALYSIS")
    add_footer(slide6, 6)

    add_card(slide6, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Hardware Specifications",
             ["Processor (CPU): Intel Core i3 / AMD Ryzen 3 or equivalent modern ARM processor.",
              "System RAM: Minimum 4 GB RAM (8 GB recommended for multi-tab browser performance).",
              "Storage Space: Less than 50 MB local storage for static HTML/JS assets and embedded training vectors.",
              "Display Resolution: 1280 x 720 minimum (1920 x 1080 recommended for optimal SVG dashboard layout).",
              "Graphics Hardware: Standard integrated graphics (No dedicated GPU required)."],
             badge="Hardware Specs")

    add_card(slide6, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Feasibility Evaluation",
             ["Low Computational Overhead: Pre-computed model centroids and matrix weights eliminate heavy client-side training.",
              "Universal Accessibility: Runs on standard laptops, desktop PCs, Chromebooks, and tablet devices.",
              "Zero GPU Dependency: All vector dot-products and Euclidean distance calculations complete in under 500ms on CPU.",
              "Conclusion: Highly feasible on existing consumer hardware without upgrades."],
             badge="Feasibility Assessment")

    # =========================================================================
    # SLIDE 7: TECHNICAL FEASIBILITY - SOFTWARE REQUIREMENTS
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, LIGHT_BG)
    add_header(slide7, "Technical Feasibility: Software Requirements", "FEASIBILITY ANALYSIS")
    add_footer(slide7, 7)

    add_card(slide7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Software Stack & Dependencies",
             ["Operating System: Windows 10/11, macOS, Linux, ChromeOS, Android/iOS.",
              "Execution Runtime: Modern Web Browser supporting ES6+ JavaScript (Chrome 90+, Firefox 88+, Safari 14+).",
              "Document Parsers: PDF.js 2.16 (PDF text extraction), Mammoth.js 1.6 (.docx XML parser).",
              "Typography & Icons: Lucide SVG icons, MathJax 3 (LaTeX rendering engine).",
              "Data Training Tools: Python 3.10+, Scikit-Learn, NLTK, NLPAug (Offline model generation)."],
             badge="Software Stack")

    add_card(slide7, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Client-Side Architecture Justification",
             ["Zero-Installation Setup: Users simply open index.html in any standard Web browser.",
              "Sandboxed Execution: Operates strictly within browser memory space without root/admin permissions.",
              "Cross-Browser Compatibility: Built using standardized W3C Web technologies (HTML5, Vanilla ES6 JS, CSS3).",
              "Conclusion: Meets all technical software feasibility benchmarks seamlessly."],
             badge="Software Feasibility")

    # =========================================================================
    # SLIDE 8: OPERATIONAL FEASIBILITY
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, LIGHT_BG)
    add_header(slide8, "Operational Feasibility Analysis", "FEASIBILITY ANALYSIS")
    add_footer(slide8, 8)

    add_card(slide8, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "User Experience (UX)",
             ["Intuitive dashboard layout modeled after modern professional analytics suites.",
              "Drag-and-drop file uploader supporting .TXT, .PDF, and .DOCX files.",
              "One-click sample text triggers for instant testing without copying files."],
             badge="01. Usability")

    add_card(slide8, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Instant Feedback Latency",
             ["Full 6-engine stylometric evaluation completes in <500ms locally.",
              "No network latency, server queuing, or downtime issues during peak academic hours.",
              "Real-time visual gauge updates immediately upon clicking 'Run Analysis'."],
             badge="02. Speed & Latency")

    add_card(slide8, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Transparent Explainability",
             ["Generates color-coded sentence heatmaps showing exact high/low perplexity zones.",
              "Interactive SVG charts explain feature distributions across 5 syntactic dimensions.",
              "Fosters trust between educators and students through complete algorithmic transparency."],
             badge="03. Transparency")

    # =========================================================================
    # SLIDE 9: ECONOMIC FEASIBILITY
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, LIGHT_BG)
    add_header(slide9, "Economic Feasibility & Cost Model", "FEASIBILITY ANALYSIS")
    add_footer(slide9, 9)

    add_card(slide9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Cost Breakdown & Savings",
             ["Server Hosting Costs: $0.00 (No backend cloud servers, database clusters, or serverless functions required).",
              "API Licensing Costs: $0.00 (No subscription fees for third-party LLM detection APIs).",
              "Open-Source Libraries: Uses open-source client libraries (PDF.js, Mammoth.js, MathJax) under MIT/Apache licenses.",
              "Maintenance Overhead: Static web assets eliminate database backups, server patch management, and security monitoring."],
             badge="Cost Structure")

    add_card(slide9, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Economic Feasibility Conclusion",
             ["Traditional cloud detectors cost thousands of dollars annually in API throughput and infrastructure.",
              "Sentinel delivers a zero-marginal-cost execution model that can be hosted on free static hosts (GitHub Pages, Vercel, Netlify).",
              "Highly cost-effective and economically viable for academic institutions and self-hosted deployments.",
              "Conclusion: 100% Economically Feasible."],
             badge="Economic Viability")

    # =========================================================================
    # SLIDE 10: SCHEDULE FEASIBILITY & ROADMAP
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, LIGHT_BG)
    add_header(slide10, "Schedule Feasibility & Sprint Roadmap", "PROJECT TIMELINE")
    add_footer(slide10, 10)

    add_card(slide10, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Agile Sprint Breakdown",
             ["Sprint 1: Background Research & Literature Review (Weeks 1-2).",
              "Sprint 2: Requirement Analysis & System Design Diagrams (Weeks 3-4).",
              "Sprint 3: Machine Learning Model Training & Feature Extraction (Weeks 5-7).",
              "Sprint 4: Client-Side JS Core Development & Module Wiring (Weeks 8-10).",
              "Sprint 5: Dashboard UI Design, SVG Graphics & MathJax Integration (Weeks 11-12).",
              "Sprint 6: Benchmark Dataset Testing, Accuracy Evaluation & Final Documentation (Weeks 13-14)."],
             badge="Sprint Breakdown")

    add_image_slide(slide10, 'extracted_images/image_03.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure: Project Development Schedule & Sprint Roadmap")

    # =========================================================================
    # SLIDE 11: DEVELOPMENT METHODOLOGY
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, LIGHT_BG)
    add_header(slide11, "Development Methodology: Agile SDLC", "METHODOLOGY")
    add_footer(slide11, 11)

    add_card(slide11, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "1. Sprint Planning & Iteration",
             ["Adopted Agile Software Development Life Cycle (SDLC).",
              "2-week sprint cycles focusing on distinct feature increments.",
              "Regular progress reviews and continuous architectural refinement."],
             badge="Agile Framework")

    add_card(slide11, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "2. Modular Component Design",
             ["Decoupled JavaScript architecture following object-oriented single-responsibility principles.",
              "Independent development and unit testing of feature extraction engines before UI integration."],
             badge="Modular Design")

    add_card(slide11, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "3. Empirical Evaluation Cycle",
             ["Continuous evaluation against a benchmark dataset of 1,000 human & AI documents.",
              "Iterative tuning of ensemble classifier weights based on confusion matrix feedback."],
             badge="Data-Driven Tuning")

    # =========================================================================
    # SLIDE 12: SYSTEM ARCHITECTURE OVERVIEW
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, LIGHT_BG)
    add_header(slide12, "System Architecture Overview", "CHAPTER 3: SYSTEM DESIGN")
    add_footer(slide12, 12)

    add_card(slide12, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0),
             "Client-Side Multi-Engine Architecture",
             ["1. Input Layer: Drag & Drop Dropzone, Text Area Editor, Sample Loaders, FileParser (.txt, .pdf, .docx).",
              "2. Pre-processing Pipeline: Sentence Splitting, Word Tokenization, Lexical Normalization, TTR Calculation.",
              "3. Analytical Engine Layer (6 Parallel Classifiers):",
              "    • Statistical Perplexity & Burstiness Engine (5% Weight)",
              "    • K-Nearest Neighbors (KNN) Feature Vector Classifier (20% Weight)",
              "    • Cosine Similarity Vector Angle Classifier (10% Weight)",
              "    • Stylometric Logistic Regression Classifier (25% Weight)",
              "    • Character Trigram Frequency Similarity Engine (25% Weight)",
              "    • POS Syntax Ratios Similarity Engine (15% Weight)",
              "4. Orchestrator & Calibration Layer: Weighted Ensemble Integration Engine + Heuristic Grammar Perfection Calibration.",
              "5. Visualization & Presentation Layer: SVG Semi-Circular Gauge, 2D KNN Scatter Plot, 5-Axis Radar Web, Sentence Heatmap Visualizer, MathJax Resolution Panel."],
             badge="End-to-End Pipeline")

    # =========================================================================
    # SLIDE 13: CORE ALGORITHMS & MATHEMATICAL MODELS
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, LIGHT_BG)
    add_header(slide13, "Core Algorithms & Mathematical Formulations", "SYSTEM ALGORITHMS")
    add_footer(slide13, 13)

    add_card(slide13, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.4),
             "1. Statistical Perplexity & Burstiness",
             ["Perplexity (PP): PP = exp( - (1/N) * sum( ln P(w_i) ) )",
              "Burstiness (B): B = ( std_dev - mean ) / ( std_dev + mean )",
              "Human text exhibits high perplexity spikes and high burstiness variance."],
             badge="Statistical NLP")

    add_card(slide13, Inches(6.8), Inches(1.6), Inches(5.6), Inches(2.4),
             "2. Stylometric Logistic Regression",
             ["Log-odds formulation: z = beta_0 + sum( beta_i * x_i )",
              "Sigmoid probability mapping: P(AI) = 1 / ( 1 + exp( -z ) )",
              "Trained on 5 normalized stylometric feature vectors."],
             badge="Machine Learning")

    add_card(slide13, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.4),
             "3. Cosine Similarity Vector Space",
             ["Vector angle calculation: cos(theta) = ( A . B ) / ( ||A|| * ||B|| )",
              "Evaluates document character trigrams and POS syntax distributions against human/AI prototypes."],
             badge="Vector Space")

    add_card(slide13, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.4),
             "4. Ensemble Score & Grammar Calibration",
             ["Base Weighted Score = sum( w_i * S_i )",
              "Final AI % = min( 100%, max( 0%, Base_Score * Grammar_Factor + Boost ) )",
              "Grammar perfection score acts as a scaling modifier for synthetic texts."],
             badge="Ensemble Fusion")

    # =========================================================================
    # SLIDE 14: CLASS DIAGRAM
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, LIGHT_BG)
    add_header(slide14, "Object Modeling: Class Diagram", "SYSTEM DESIGN")
    add_footer(slide14, 14)

    add_card(slide14, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Object Architecture Details",
             ["app.js Controller: Coordinates events, UI state, and calls classification modules.",
              "FileParser Class: Wraps PDF.js and Mammoth.js to extract plain text.",
              "PerplexityAnalyzer Class: Computes n-gram log-likelihoods and burstiness.",
              "KNNClassifier Class: Holds training centroid coordinates and computes Euclidean distances.",
              "CosineSimilarityAnalyzer Class: Executes dot-product angle calculations.",
              "ChartRenderer Class: Generates SVG paths for Gauge, Scatter Plot, and Radar Web."],
             badge="Class Structure")

    add_image_slide(slide14, 'extracted_images/image_06.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure 3.3: Class Diagram for AI Text Detector")

    # =========================================================================
    # SLIDE 15: STATE DIAGRAM
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, LIGHT_BG)
    add_header(slide15, "Dynamic Modeling: State Diagram", "DYNAMIC MODELING")
    add_footer(slide15, 15)

    add_card(slide15, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "State Transition Lifecycle",
             ["1. IDLE State: Waiting for user text input or file drop.",
              "2. DOCUMENT LOADED State: Text extracted and word count validated (>= 10 words).",
              "3. TOKENIZING State: Text split into sentences and words; TTR/burstiness computed.",
              "4. ANALYZING State: 6 base classifier engines execute in parallel.",
              "5. SCORING State: Weighted ensemble score and grammar factor calculated.",
              "6. RENDERING State: Dashboard charts, sentence heatmaps, and LaTeX panel updated."],
             badge="State Transitions")

    add_image_slide(slide15, 'extracted_images/image_07.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure 3.4: State Diagram for AI Text Detector")

    # =========================================================================
    # SLIDE 16: SEQUENCE DIAGRAM
    # =========================================================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16, LIGHT_BG)
    add_header(slide16, "Dynamic Modeling: Sequence Diagram", "DYNAMIC MODELING")
    add_footer(slide16, 16)

    add_card(slide16, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Asynchronous Event Sequence",
             ["User actions trigger event listeners in app.js.",
              "1. User drops file -> FileParser parses text asynchronously.",
              "2. app.js validates text length -> returns error if < 10 words.",
              "3. app.js passes tokenized text to Perplexity, KNN, Cosine, Logistic, and Grammar analyzers.",
              "4. Analyzers return metric percentages to app.js.",
              "5. app.js computes ensemble score and passes data to ChartRenderer.",
              "6. ChartRenderer updates DOM SVG paths and invokes MathJax rendering."],
             badge="Message Control Flow")

    add_image_slide(slide16, 'extracted_images/image_08.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure 3.5: Sequence Diagram for AI Text Detector")

    # =========================================================================
    # SLIDE 17: ACTIVITY DIAGRAM
    # =========================================================================
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17, LIGHT_BG)
    add_header(slide17, "Process Modeling: Activity Diagram", "PROCESS MODELING")
    add_footer(slide17, 17)

    add_card(slide17, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Operational Workflow",
             ["Process starts upon text submission.",
              "Decision Node: Is word count >= 10?",
              "• NO: Display warning prompt and abort execution.",
              "• YES: Proceed with local sentence tokenization.",
              "Fork Node: Parallel feature extraction across 6 engines.",
              "Join Node: Combine outputs into weighted ensemble equation.",
              "Apply Heuristic Grammar Calibration multiplier.",
              "Render Gauge, Scatter Plot, Radar Web, Heatmap highlights, and MathJax formulas."],
             badge="Workflow Decision Nodes")

    add_image_slide(slide17, 'extracted_images/image_09.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure 3.6: Activity Diagram for AI Text Detector")

    # =========================================================================
    # SLIDE 18: COMPONENT DIAGRAM
    # =========================================================================
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18, LIGHT_BG)
    add_header(slide18, "System Design: Component Diagram", "SYSTEM DESIGN")
    add_footer(slide18, 18)

    add_card(slide18, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Component Interfaces",
             ["HTML5 Dashboard: Front-end UI view containing input dropzone and output grid.",
              "app.js Controller: Core orchestrator linking UI components to analytical engines.",
              "FileParser Component: Integrates PDF.js & Mammoth.js third-party libraries.",
              "Analysis Engines: Standalone modular classes for Perplexity, KNN, Cosine, Logistic Regression.",
              "ChartRenderer Component: Interfaces with SVG DOM graphic elements.",
              "MathJax Component: Typesets LaTeX math equations on demand."],
             badge="Component Architecture")

    add_image_slide(slide18, 'extracted_images/image_10.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure 3.7: Component Diagram for AI Text Detector")

    # =========================================================================
    # SLIDE 19: DEPLOYMENT DIAGRAM
    # =========================================================================
    slide19 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide19, LIGHT_BG)
    add_header(slide19, "System Design: Deployment Diagram", "DEPLOYMENT MODEL")
    add_footer(slide19, 19)

    add_card(slide19, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Local Client Execution Node",
             ["Client Device Node: User's PC / Laptop / Mobile device.",
              "Browser Sandbox Artifact: Web Browser Runtime (V8 / SpiderMonkey engine).",
              "Static Web Files: index.html, style.css, app.js stored in browser cache.",
              "In-Memory Data Structures: Pre-computed model weights and feature vectors.",
              "FileReader API: Local file system access without remote server uploads.",
              "Guarantees 100% offline security and complete data isolation."],
             badge="Deployment Node")

    add_image_slide(slide19, 'extracted_images/image_11.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure 3.8: Deployment Diagram for AI Text Detector")

    # =========================================================================
    # SLIDE 20: TOOLS & TECHNOLOGIES USED
    # =========================================================================
    slide20 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide20, LIGHT_BG)
    add_header(slide20, "Tools & Technologies Stack", "CHAPTER 4: IMPLEMENTATION")
    add_footer(slide20, 20)

    add_card(slide20, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Front-End Web Stack",
             ["HTML5: Semantic dashboard layout structure.",
              "CSS3: Modern dark glassmorphic design, responsive grid layouts, custom variables.",
              "Vanilla ES6+ JavaScript: Modular object-oriented execution engine without bulky framework overhead."],
             badge="Core Web Tools")

    add_card(slide20, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Client-Side Libraries",
             ["PDF.js (v2.16): Mozilla's canvas PDF text stream parser.",
              "Mammoth.js (v1.6): Word .docx raw XML text converter.",
              "MathJax 3: LaTeX typesetting for mathematical formulas.",
              "Lucide Icons: Modern clean vector dashboard icon set."],
             badge="External Libraries")

    add_card(slide20, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "Data Science & Training",
             ["Python 3.10: Data preparation and offline model training.",
              "Scikit-Learn: Logistic Regression & KNN model generation.",
              "NLTK: Corpus tokenization and POS tagging.",
              "NLPAug: Adversarial data augmentation & evasion testing."],
             badge="Python Data Stack")

    # =========================================================================
    # SLIDE 21: SYSTEM MODULES OVERVIEW
    # =========================================================================
    slide21 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide21, LIGHT_BG)
    add_header(slide21, "System Implementation Modules Overview", "SYSTEM MODULES")
    add_footer(slide21, 21)

    add_card(slide21, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Core Analytical Engine Modules",
             ["Module 1: FileParser Module (.txt, .pdf, .docx document parser).",
              "Module 2: Perplexity & Burstiness Engine (N-gram likelihood analyzer).",
              "Module 3: KNN Feature Vector Classifier (2D spatial Euclidean distance).",
              "Module 4: Stylometric Logistic Regression Classifier (5D feature model).",
              "Module 5: Character Trigram Engine (Sub-word frequency similarity).",
              "Module 6: POS Syntax Ratios Engine (Part-of-speech distribution analyzer).",
              "Module 7: Heuristic Grammar & Typography Calibrator."],
             badge="Analytical Engines")

    add_card(slide21, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "UI, Controller & Visualization Modules",
             ["Module 8: Ensemble Scoring & Grammar Calibration Orchestrator.",
              "Module 9: ChartRenderer Gauge Module (SVG semi-circular dial).",
              "Module 10: ChartRenderer Scatter Plot & Radar Web Visualizer.",
              "Module 11: Sentence-Level Perplexity Heatmap Engine.",
              "Module 12: MathJax Algorithmic Resolution Panel.",
              "Module 13: app.js Controller & Event Listener Sandbox."],
             badge="Presentation & UI")

    # =========================================================================
    # SLIDE 22: MODULE 1 - FILEPARSER MODULE
    # =========================================================================
    slide22 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide22, LIGHT_BG)
    add_header(slide22, "Module 1: FileParser Class Implementation", "MODULE DETAILS")
    add_footer(slide22, 22)

    add_card(slide22, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "File Processing Capabilities",
             ["Supports multiple file formats seamlessly:",
              "• Plain Text (.TXT / .MD): Direct FileReader API text extraction.",
              "• PDF Documents (.PDF): Integrates PDF.js to iterate through document pages, extract text items, and normalize whitespace.",
              "• Word Documents (.DOCX): Integrates Mammoth.js to parse raw document XML and extract clean body paragraphs.",
              "Automated Error Handling: Gracefully catches corrupt files, password protection locks, and invalid formats."],
             badge="File Parsing Engine")

    add_card(slide22, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Text Pre-processing Pipeline",
             ["Sentence Tokenization: Regex-based sentence boundary detection handling abbreviation exceptions.",
              "Word Tokenization: Alphanumeric extraction with lowercasing for stylometric feature extraction.",
              "Length Validation: Enforces a 10-word minimum threshold required for statistically reliable stylometric analysis.",
              "Session Cache Storage: Caches raw text and tokens to prevent redundant parsing during analysis runs."],
             badge="Pre-processing Pipeline")

    # =========================================================================
    # SLIDE 23: MODULE 2 - PERPLEXITY & BURSTINESS ENGINE
    # =========================================================================
    slide23 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide23, LIGHT_BG)
    add_header(slide23, "Module 2: Perplexity & Burstiness Engine", "MODULE DETAILS")
    add_footer(slide23, 23)

    add_card(slide23, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Perplexity Calculation Mechanics",
             ["Evaluates log-likelihood of token sequences using a localized n-gram language model.",
              "LLMs generate highly predictable text with consistently low perplexity scores.",
              "Human text exhibits high perplexity spikes due to unexpected word choices and creative phrasing.",
              "Perplexity Score = exp( - (1/N) * sum( log( P(w_i | context) ) ) )",
              "Low Perplexity (< 40) -> Strong indicator of AI generation."],
             badge="Perplexity Metric")

    add_card(slide23, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Burstiness Calculation Mechanics",
             ["Measures variation in sentence lengths across the document.",
              "Human writers naturally alternate between short punchy sentences and complex compound structures.",
              "LLMs produce uniform sentence lengths with low variance.",
              "Burstiness B = ( std_dev - mean ) / ( std_dev + mean )",
              "Bounded in range [-1, +1]. High burstiness indicates authentic human authorship."],
             badge="Burstiness Metric")

    # =========================================================================
    # SLIDE 24: MODULE 3 - KNN CLASSIFIER MODULE
    # =========================================================================
    slide24 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide24, LIGHT_BG)
    add_header(slide24, "Module 3: KNN Feature Vector Classifier", "MODULE DETAILS")
    add_footer(slide24, 24)

    add_card(slide24, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "2D Feature Space Mapping",
             ["Feature Dimension 1 (X-axis): Document N-gram Perplexity Score.",
              "Feature Dimension 2 (Y-axis): Sentence Length Standard Deviation (Burstiness Variance).",
              "Pre-loaded Centroids: Embedded training set reference coordinates for Human writing clusters vs. AI text clusters.",
              "Normalizes raw features using Min-Max scaling into bounded [0, 1] range."],
             badge="Feature Space")

    add_card(slide24, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Classification Algorithm",
             ["Computes Euclidean distances to nearest K=5 training reference points:",
              "d(P, Q) = sqrt( (x_2 - x_1)^2 + (y_2 - y_1)^2 )",
              "Applies distance-weighted voting algorithm to assign local AI probability.",
              "Provides coordinate data directly to the ChartRenderer for 2D visual scatter plot plotting.",
              "Model Weight in Ensemble: 20%."],
             badge="Distance Classifier")

    # =========================================================================
    # SLIDE 25: MODULE 4 - STYLOMETRIC LOGISTIC REGRESSION
    # =========================================================================
    slide25 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide25, LIGHT_BG)
    add_header(slide25, "Module 4: Stylometric Logistic Regression", "MODULE DETAILS")
    add_footer(slide25, 25)

    add_card(slide25, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "5-Dimensional Stylometric Feature Vector",
             ["1. Lexical Diversity (Type-Token Ratio / TTR): Unique words / Total words.",
              "2. Burstiness Index: Sentence length variance measure.",
              "3. AI Vocabulary Density: Frequency of characteristic AI transition terms ('pivotal', 'crucial', 'delve', 'testament').",
              "4. Coleman-Liau Readability Index (CLI): Grade level readability measure.",
              "5. Punctuation Density: Ratio of commas and punctuation marks per 100 words."],
             badge="5D Feature Vector")

    add_card(slide25, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Model Formulation & Weighting",
             ["Logit formulation: z = -2.45 + (1.82 * V1) + (2.15 * V2) + (3.40 * V3) - (0.95 * V4) + (1.10 * V5)",
              "Probability mapping: P(AI) = 1 / ( 1 + exp( -z ) )",
              "Trained on 1,000 document samples using Scikit-Learn offline.",
              "Highest individual model precision (95.1%).",
              "Model Weight in Ensemble: 25%."],
             badge="Supervised ML Model")

    # =========================================================================
    # SLIDE 26: MODULE 5 - CHARACTER TRIGRAM ENGINE
    # =========================================================================
    slide26 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide26, LIGHT_BG)
    add_header(slide26, "Module 5: Character Trigram Engine", "MODULE DETAILS")
    add_footer(slide26, 26)

    add_card(slide26, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Character 3-Gram Extraction",
             ["Extracts overlapping character triplets across the text (e.g. 'the' -> ['the', 'he ']).",
              "Captures subtle sub-word morpho-syntactic habits independent of high-level vocabulary.",
              "Generates a 20-dimensional character trigram frequency vector V_trigram.",
              "Robust against synonym substitution and simple word-level paraphrasing evasions."],
             badge="Sub-Word Extraction")

    add_card(slide26, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Cosine Angle Evaluation",
             ["Compares document trigram vector against reference AI and Human prototype vectors:",
              "Sim(Doc, AI) = ( V_doc . V_ai ) / ( ||V_doc|| * ||V_ai|| )",
              "Converts angular similarity into a normalized AI percentage score.",
              "Extremely effective for detecting machine-generated formatting artifacts.",
              "Model Weight in Ensemble: 25%."],
             badge="Vector Similarity")

    # =========================================================================
    # SLIDE 27: MODULE 6 - POS SYNTAX RATIOS ENGINE
    # =========================================================================
    slide27 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide27, LIGHT_BG)
    add_header(slide27, "Module 6: POS Syntax Ratios Classifier", "MODULE DETAILS")
    add_footer(slide27, 27)

    add_card(slide27, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Part-of-Speech Structural Profiling",
             ["Analyzes grammatical structure across 5 syntactic dimensions:",
              "• Noun Ratio: Density of nouns relative to total token count.",
              "• Verb Ratio: Action verb frequency distribution.",
              "• Adjective / Adverb Density: Modifying word usage ratios.",
              "• Function Word Ratio: Prepositions, conjunctions, and articles.",
              "• Passive Voice & Clause Density: Complex syntactic subordination."],
             badge="Syntactic Profiling")

    add_card(slide27, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Syntactic Prototype Matching",
             ["LLMs exhibit highly rigid syntactic distribution ratios across paragraphs.",
              "Human writing demonstrates dynamic structural shifts based on rhetorical tone.",
              "Measures 5-axis Euclidean distance to human vs. AI syntactic profiles.",
              "Provides 5-axis coordinate data to the SVG Radar Web visualizer.",
              "Model Weight in Ensemble: 15%."],
             badge="Structural Matching")

    # =========================================================================
    # SLIDE 28: MODULE 7 - HEURISTIC GRAMMAR & TYPOGRAPHY
    # =========================================================================
    slide28 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide28, LIGHT_BG)
    add_header(slide28, "Module 7: Heuristic Grammar & Typography", "MODULE DETAILS")
    add_footer(slide28, 28)

    add_card(slide28, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Grammar Perfection Evaluator",
             ["Scans text for common human writing flaws: Subject-verb disagreement, run-on sentences, missing commas, double spaces, capitalization errors.",
              "Generates a normalized Grammar Perfection Score (0% to 100%).",
              "LLM outputs almost universally exhibit 100% flawless grammar and perfect punctuation alignment."],
             badge="Grammar Evaluation")

    add_card(slide28, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Ensemble Calibration Modifier",
             ["Calculates Grammar Factor = Perfection Score / 100.",
              "Applies Perfect Grammar Boost (+15.0%) when Perfection Score = 100%.",
              "Applies Grammar Penalty discount when human grammatical typos are detected.",
              "Prevents false positives on human text containing natural typing errors."],
             badge="Calibration Logic")

    # =========================================================================
    # SLIDE 29: MODULE 8 - ENSEMBLE SCORING ORCHESTRATOR
    # =========================================================================
    slide29 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide29, LIGHT_BG)
    add_header(slide29, "Module 8: Ensemble Scoring Orchestrator", "ENSEMBLE FUSION")
    add_footer(slide29, 29)

    add_card(slide29, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0),
             "Unified Ensemble Fusion Pipeline",
             ["1. Base Weighted Score Calculation:",
              "    Base Score = (Logistic * 0.25) + (Trigram * 0.25) + (KNN * 0.20) + (POS * 0.15) + (Cosine * 0.10) + (Perplexity * 0.05)",
              "2. Grammar Calibration Adjustment:",
              "    Final AI % = min( 100.0%, max( 0.0%, ( Base Score * Grammar Factor ) + Perfect Grammar Boost ) )",
              "3. Classification Threshold Mapping:",
              "    • 0.0% - 35.0%   -> Negligible / Human-Like Text (Sage Green)",
              "    • 35.1% - 50.0%  -> Low AI Probability (Blue-Green)",
              "    • 50.1% - 60.0%  -> Borderline / Uncertain (Amber Yellow)",
              "    • 60.1% - 80.0%  -> High AI Probability (Orange)",
              "    • 80.1% - 100.0% -> Extremely High AI Probability (Terracotta Red)"],
             badge="Weighted Voting Engine")

    # =========================================================================
    # SLIDE 30: MODULE 9 - SVG GAUGE & CHART RENDERER
    # =========================================================================
    slide30 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide30, LIGHT_BG)
    add_header(slide30, "Module 9: SVG Semi-Circular Gauge Renderer", "VISUALIZATION ENGINE")
    add_footer(slide30, 30)

    add_card(slide30, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Dynamic SVG Dial Rendering",
             ["Generates clean resolution-independent SVG vector graphics dynamically in the DOM.",
              "Calculates smooth arc trajectories using parametric trigonometric equations:",
              "x = cx + r * cos(angle),  y = cy + r * sin(angle)",
              "Dynamic Color Interpolation: Seamlessly transitions color from Sage Green (#22C55E) -> Yellow (#EAB308) -> Red (#EF4444).",
              "Displays bold central percentage score and category badge label."],
             badge="SVG Mathematics")

    add_card(slide30, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "UI Integration & Animation",
             ["CSS Transitions: Smooth 800ms ease-out needle rotation animation upon analysis completion.",
              "High Contrast Glassmorphic Container: Styled with dark slate background and glowing border highlights.",
              "Fully responsive across all screen sizes and mobile viewports.",
              "Zero reliance on heavy external charting libraries (Canvas / Chart.js)."],
             badge="UI Presentation")

    # =========================================================================
    # SLIDE 31: MODULE 10 - SCATTER PLOT & RADAR WEB
    # =========================================================================
    slide31 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide31, LIGHT_BG)
    add_header(slide31, "Module 10: 2D KNN Scatter & 5-Axis Radar Web", "VISUALIZATION ENGINE")
    add_footer(slide31, 31)

    add_card(slide31, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "2D KNN Feature Scatter Plot",
             ["Plots 2D feature coordinates (Perplexity vs. Sentence Variance).",
              "Renders reference training clusters: Human benchmark dots (Green) vs. AI benchmark dots (Red).",
              "Plots a distinct pulsing coordinate marker for the active document.",
              "Visually demonstrates spatial clustering proximity to the user."],
             badge="2D Scatter Plot")

    add_card(slide31, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "5-Axis Syntactic Radar Web",
             ["Renders a 5-axis spider web polygon (Noun, Verb, Adj, Function Word, Clause density).",
              "Overlays document syntactic profile against Human and AI prototype boundaries.",
              "Highlights exact syntactic deviations responsible for the classification decision.",
              "Dramatically improves model explainability."],
             badge="5-Axis Radar Web")

    # =========================================================================
    # SLIDE 32: MODULE 11 - SENTENCE PERPLEXITY HEATMAP
    # =========================================================================
    slide32 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide32, LIGHT_BG)
    add_header(slide32, "Module 11: Sentence Perplexity Heatmap", "VISUAL EXPLAINABILITY")
    add_footer(slide32, 32)

    add_card(slide32, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Sentence-Level Granular Analysis",
             ["Splits target document into individual sentence tokens.",
              "Computes localized n-gram perplexity score for each isolated sentence.",
              "Maps sentence perplexity onto a color gradient scale:",
              "• Low Perplexity (< 35): Red / Orange background highlight (High AI likelihood).",
              "• Moderate Perplexity (35-65): Yellow background highlight.",
              "• High Perplexity (> 65): Green / Clear background highlight (Authentic human)."],
             badge="Perplexity Mapping")

    add_card(slide32, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Interactive Editor Integration",
             ["Renders color highlights directly inside the interactive text display sandbox.",
              "Tooltip Hover Feature: Displays exact sentence word count and numerical perplexity value when hovering over any sentence.",
              "Allows users and educators to instantly locate AI-generated sections within hybrid documents.",
              "Key feature for academic integrity verification."],
             badge="Editor Integration")

    # =========================================================================
    # SLIDE 33: MODULE 12 - MATHJAX RESOLUTION PANEL
    # =========================================================================
    slide33 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide33, LIGHT_BG)
    add_header(slide33, "Module 12: MathJax Resolution Panel", "TRANSPARENT DIAGNOSTICS")
    add_footer(slide33, 33)

    add_card(slide33, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Dynamic LaTeX Formula Generation",
             ["Generates exact LaTeX mathematical strings dynamically upon analysis completion.",
              "Formats 5D feature vector coordinates V = [v1, v2, v3, v4, v5].",
              "Displays individual model scores: S_logistic, S_trigram, S_knn, S_pos, S_cosine, S_perplexity.",
              "Typesets weighted sum equation and grammar calibration modifier step-by-step."],
             badge="LaTeX Engine")

    add_card(slide33, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "MathJax 3 Typesetting Engine",
             ["Invokes MathJax.typesetPromise() to render publication-grade math formulas in SVG format.",
              "Provides full auditability for researchers, peer reviewers, and defense committees.",
              "Proves that Sentinel operates on transparent mathematical principles rather than arbitrary heuristics.",
              "Eliminates black-box opacity entirely."],
             badge="Typesetting Integration")

    # =========================================================================
    # SLIDE 34: MODULE 13 - EVENT CONTROLLER & PRIVACY
    # =========================================================================
    slide34 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide34, LIGHT_BG)
    add_header(slide34, "Module 13: Event Controller & Privacy Sandbox", "CLIENT SANDBOX")
    add_footer(slide34, 34)

    add_card(slide34, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Event-Driven JavaScript Architecture",
             ["Listens for DOM events: Dragover, Drop, File Select, Text Input, Sample Click, Clear Button.",
              "Debounces text input events to prevent UI lag during typing.",
              "Manages UI state toggles: Loading progress bar, active tab highlights, error prompts.",
              "Simulates asynchronous progress steps (0% -> 100%) for smooth visual UX."],
             badge="Event Controller")

    add_card(slide34, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "100% Privacy Sandbox Guarantee",
             ["All text processing occurs in browser volatile memory (RAM).",
              "Zero HTTP requests or API calls are dispatched during document analysis.",
              "No localStorage or IndexedDB persistence of user document text.",
              "Closing the browser tab immediately purges all analyzed text from memory.",
              "Compliant with GDPR, FERPA, and strict organizational privacy policies."],
             badge="Privacy Isolation")

    # =========================================================================
    # SLIDE 35: SYSTEM TESTING STRATEGY
    # =========================================================================
    slide35 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide35, LIGHT_BG)
    add_header(slide35, "System Testing & Verification Strategy", "CHAPTER 4: TESTING")
    add_footer(slide35, 35)

    add_card(slide35, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "1. Unit Testing",
             ["Verified individual JavaScript classes independently.",
              "Tested FileParser with corrupted PDF/DOCX streams.",
              "Validated mathematical edge cases (e.g. 0-variance text, single-sentence inputs)."],
             badge="Unit Tests")

    add_card(slide35, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "2. Integration Testing",
             ["Verified end-to-end data flow between FileParser, Controller, Classifiers, and ChartRenderer.",
              "Ensured SVG graphics re-render correctly across window resize events."],
             badge="Integration Tests")

    add_card(slide35, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.0),
             "3. Evasion Robustness Testing",
             ["Subjected detector to adversarial perturbations using NLPAug library:",
              "• Synonym substitution",
              "• Typo insertion",
              "• Paraphrasing & sentence reordering",
              "Evaluated detector robustness under evasion attempts."],
             badge="Adversarial Testing")

    # =========================================================================
    # SLIDE 36: CLASSIFIER PERFORMANCE RESULTS
    # =========================================================================
    slide36 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide36, LIGHT_BG)
    add_header(slide36, "Classifier Performance Benchmark Results", "RESULT ANALYSIS")
    add_footer(slide36, 36)

    # Performance Metrics Table Shape
    rows = 8
    cols = 5
    table_shape = slide36.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(6.5), Inches(5.0))
    table = table_shape.table
    
    headers = ["Model Engine", "Accuracy", "Precision", "Recall", "F1 Score"]
    data = [
        ["Stylometric Logistic Reg.", "93.4%", "95.1%", "91.8%", "93.4%"],
        ["Character Trigram Cosine", "91.8%", "93.2%", "90.1%", "91.6%"],
        ["K-Nearest Neighbors (KNN)", "89.5%", "90.8%", "87.9%", "89.3%"],
        ["POS Syntax Ratios", "86.2%", "87.5%", "84.6%", "86.0%"],
        ["Cosine Vector Analyzer", "84.1%", "85.0%", "82.8%", "83.9%"],
        ["Perplexity & Burstiness", "81.5%", "82.4%", "80.2%", "81.3%"],
        ["UNIFIED ENSEMBLE (Ours)", "94.2%", "94.8%", "93.5%", "94.1%"]
    ]
    
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_LIGHT
            
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r+1, c)
            cell.text = val
            cell.fill.solid()
            if r == 6: # Highlight ensemble
                cell.fill.fore_color.rgb = RGBColor(219, 234, 254)
            else:
                cell.fill.fore_color.rgb = CARD_BG if r%2==0 else RGBColor(241, 245, 249)
            for p in cell.text_frame.paragraphs:
                if r == 6 or c == 0:
                    p.font.bold = True
                p.font.size = Pt(11)
                p.font.color.rgb = PRIMARY_BLUE if r==6 else TEXT_DARK

    add_image_slide(slide36, 'extracted_images/image_13.png', Inches(7.5), Inches(1.6), Inches(5.0), Inches(5.0),
                    "Figure 4.19: Classifier Performance Comparison Chart")

    # =========================================================================
    # SLIDE 37: DATASET SCORE DISTRIBUTION & TRENDS
    # =========================================================================
    slide37 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide37, LIGHT_BG)
    add_header(slide37, "Dataset Distribution & Category Trends", "RESULT ANALYSIS")
    add_footer(slide37, 37)

    add_card(slide37, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Benchmark Dataset Analysis",
             ["Evaluated on 1,000 document samples across 12 writing domains.",
              "Distribution Findings:",
              "• Negligible / Human-Like: 337 docs (33.7%)",
              "• Extremely High AI: 332 docs (33.2%)",
              "• High AI Probability: 159 docs (15.9%)",
              "• Low AI Probability: 99 docs (9.9%)",
              "• Borderline / Uncertain: 73 docs (7.3%)",
              "Strong bimodal distribution confirms distinct feature separation between human and synthetic texts."],
             badge="Distribution Summary")

    add_image_slide(slide37, 'extracted_images/image_12.png', Inches(5.2), Inches(1.6), Inches(3.6), Inches(5.0),
                    "Figure 4.18: AI Probability Distribution")

    add_image_slide(slide37, 'extracted_images/image_14.png', Inches(9.0), Inches(1.6), Inches(3.5), Inches(5.0),
                    "Figure 4.20: Trend by Document Category")

    # =========================================================================
    # SLIDE 38: CONFUSION MATRIX & SAMPLE CALCULATION
    # =========================================================================
    slide38 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide38, LIGHT_BG)
    add_header(slide38, "Confusion Matrix & Step-by-Step Walkthrough", "MODEL EVALUATION")
    add_footer(slide38, 38)

    add_card(slide38, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Sample Text Evaluation Case Study",
             ["Sample Input: 103 words across 5 sentences.",
              "1. Feature Extraction: V = [1.000, 0.752, 1.000, 0.988, 0.728]",
              "2. Engine Scores:",
              "    • Logistic Reg: 97.67% | Trigram: 55.12%",
              "    • KNN: 100.00%        | POS: 50.00%",
              "    • Cosine: 100.00%     | Perplexity: 100.00%",
              "3. Base Weighted Score: 78.67%",
              "4. Grammar Calibration: 100% perfect grammar -> +15.0% boost.",
              "5. Final AI Score: min(100%, 78.67% * 1.0 + 15%) = 90.46% (Extremely High AI)."],
             badge="Case Study Walkthrough")

    add_image_slide(slide38, 'extracted_images/image_16.png', Inches(5.3), Inches(1.6), Inches(7.2), Inches(4.8),
                    "Figure 4.22: Confusion Matrix for AI Text Detector")

    # =========================================================================
    # SLIDE 39: CONCLUSION & FUTURE RECOMMENDATIONS
    # =========================================================================
    slide39 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide39, LIGHT_BG)
    add_header(slide39, "Conclusion & Future Recommendations", "CHAPTER 5: CONCLUSION")
    add_footer(slide39, 39)

    add_card(slide39, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Project Achievements",
             ["Successfully built a 100% client-side AI text detector ('Sentinel').",
              "Achieved 94.2% overall detection accuracy on 1,000 benchmark documents.",
              "Completely eliminates cloud data privacy risks and subscription API costs.",
              "Delivers sub-500ms execution latency with rich interactive SVG charts and sentence-level heatmaps."],
             badge="Key Achievements")

    add_card(slide39, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0),
             "Future Work & Roadmap",
             ["1. WebAssembly ONNX Runtime: Integrate quantized small LLM models (e.g. MobileBERT) via WASM.",
              "2. Browser Extension: Develop Chrome/Firefox extension for inline text field scanning.",
              "3. Multi-Lingual Support: Expand character trigram models for Spanish, French, and Nepali text.",
              "4. PDF Scan Report Export: Allow users to download signed verification certificates."],
             badge="Future Recommendations")

    # =========================================================================
    # SLIDE 40: SYSTEM SHOWCASE & DEMO (Dark Executive Theme)
    # =========================================================================
    slide40 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide40, DARK_BG)
    add_header(slide40, "Interactive System Showcase & Live Demonstration", "LIVE DEMO", dark=True)
    add_footer(slide40, 40, dark=True)

    add_card(slide40, Inches(0.8), Inches(1.6), Inches(4.2), Inches(5.0),
             "Sentinel Dashboard Summary",
             ["Key System Highlights:",
              "• Drag & Drop File Parser (.TXT, .PDF, .DOCX)",
              "• Real-Time Semi-Circular Gauge Score Dial",
              "• 2D KNN Feature Space Scatter Plot",
              "• 5-Axis Syntactic Radar Web Visualizer",
              "• Interactive Sentence Perplexity Heatmap Editor",
              "• Dynamic MathJax LaTeX Calculation Panel",
              "",
              "Thank You! Questions & Answers."],
             badge="Live Demo Invitation", bg_color=SLATE_CARD, border_color=PRIMARY_BLUE, dark_text=False)

    add_image_slide(slide40, 'extracted_images/image_17.png', Inches(5.3), Inches(1.6), Inches(3.5), Inches(5.0),
                    "Sentinel Dashboard Overview")

    add_image_slide(slide40, 'extracted_images/image_20.png', Inches(9.0), Inches(1.6), Inches(3.5), Inches(5.0),
                    "Gauge & Diagnostic Charts")

    # Save presentation
    output_path = "AI_Text_Detector_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
